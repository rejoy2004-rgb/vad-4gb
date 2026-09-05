"""Build the 2-slide submission deck (and its charts) from the measured runs."""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).parent
BG, PANEL = RGBColor(0x10, 0x13, 0x1A), RGBColor(0x1A, 0x1F, 0x2B)
TEXT, MUTED = RGBColor(0xE8, 0xEC, 0xF4), RGBColor(0x8B, 0x93, 0xA7)
ACCENT, GREEN, RED = RGBColor(0xF5, 0xA6, 0x23), RGBColor(0x4A, 0xDE, 0x80), RGBColor(0xF8, 0x71, 0x71)

# measured on the 34-video public test set
STAGES = [
    ("Window head,\nabsolute thresholds", 0.6875, 0.5333, 0.2000, 0.4736),
    ("+ per-video\nrelative scoring", 0.6875, 0.6576, 0.2960, 0.5471),
    ("+ clip head\nfor Level 1", 0.7292, 0.6576, 0.2960, 0.5609),
    ("+ clip head classifies\nspans, wider merge", 0.7292, 0.6737, 0.4145, 0.6058),
    ("+ 3-seed ensemble,\nmatched tuning", 0.7500, 0.6737, 0.4301, 0.6179),
]


def _hex(c):
    return "#%02X%02X%02X" % (c[0], c[1], c[2])


def chart_progression(path):
    fig, ax = plt.subplots(figsize=(7.4, 3.05), dpi=220)
    fig.patch.set_facecolor(_hex(PANEL))
    ax.set_facecolor(_hex(PANEL))
    xs = range(len(STAGES))
    ov = [s[4] for s in STAGES]
    ax.plot(xs, ov, color=_hex(ACCENT), lw=3, marker="o", ms=8,
            mfc=_hex(ACCENT), mec=_hex(PANEL), mew=2, zorder=3, label="Overall")
    for lbl, idx, col in (("Level 1", 1, "#7DD3FC"), ("Level 2", 2, "#A78BFA"),
                          ("Level 3", 3, "#4ADE80")):
        ax.plot(xs, [s[idx] for s in STAGES], lw=1.8, marker="o", ms=4.5,
                color=col, alpha=.85, label=lbl)
    # labels sit below the overall line: above it they collide with Level 2
    for x, v in zip(xs, ov):
        ax.annotate(f"{v:.3f}", (x, v), textcoords="offset points", xytext=(0, -17),
                    ha="center", color=_hex(ACCENT), fontsize=10, fontweight="bold")
    ax.axhline(0.2931, color=_hex(RED), ls="--", lw=1.3, alpha=.8)
    ax.annotate("0.293  \"call everything anomalous\"", (len(STAGES) - 1, 0.2931),
                xytext=(-6, -15), textcoords="offset points", ha="right",
                color=_hex(RED), fontsize=8.5, alpha=.95)
    ax.set_xticks(list(xs))
    ax.set_xticklabels([s[0] for s in STAGES], color=_hex(MUTED), fontsize=8.4)
    ax.set_xlim(-0.35, len(STAGES) - 0.65)
    ax.set_ylim(0.15, 0.83)
    ax.tick_params(colors=_hex(MUTED), labelsize=8.5, length=0)
    for s in ax.spines.values():
        s.set_visible(False)
    ax.grid(axis="y", color="#2C3444", lw=.8)
    ax.set_axisbelow(True)
    leg = ax.legend(loc="upper left", frameon=False, fontsize=8.6, ncol=4,
                    handlelength=1.4, columnspacing=1.1)
    for t in leg.get_texts():
        t.set_color(_hex(MUTED))
    fig.tight_layout(pad=0.5)
    fig.savefig(path, facecolor=fig.get_facecolor())
    plt.close(fig)


def chart_latency(path):
    fig, ax = plt.subplots(figsize=(3.5, 2.35), dpi=220)
    fig.patch.set_facecolor(_hex(PANEL))
    ax.set_facecolor(_hex(PANEL))
    names = ["Real time\nbudget", "Ours\n(end to end)"]
    vals = [1.0, 20.3]
    bars = ax.barh(names, vals, color=[_hex(MUTED), _hex(ACCENT)], height=.5)
    ax.bar_label(bars, labels=["1.0x", "20.3x"], color=_hex(TEXT),
                 fontsize=11, fontweight="bold", padding=5)
    ax.set_xlim(0, 30)   # headroom so the value label is not clipped
    ax.set_xticks([])
    ax.tick_params(colors=_hex(MUTED), labelsize=9.5, length=0)
    for s in ax.spines.values():
        s.set_visible(False)
    fig.tight_layout(pad=0.4)
    fig.savefig(path, facecolor=fig.get_facecolor())
    plt.close(fig)


# ------------------------------------------------------------------ pptx utils
def box(slide, x, y, w, h, fill=PANEL, line=None, radius=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    if radius:
        shp.adjustments[0] = radius
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, size=12, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, {})]
    for i, (txt, style) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = style.get("align", align)
        p.line_spacing = style.get("spacing", spacing)
        if style.get("space_before"):
            p.space_before = Pt(style["space_before"])
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(style.get("size", size))
        r.font.bold = style.get("bold", bold)
        r.font.color.rgb = style.get("color", color)
        r.font.name = style.get("font", "Segoe UI")
    return tb


def stat(slide, x, y, w, value, label, color=ACCENT):
    box(slide, x, y, w, 1.06, PANEL, radius=0.12)
    text(slide, x + 0.16, y + 0.13, w - 0.32, 0.45, value, size=23, bold=True, color=color)
    text(slide, x + 0.16, y + 0.62, w - 0.32, 0.34, label, size=9.5, color=MUTED)


def chip(slide, x, y, w, h, label, color):
    box(slide, x, y, w, h, PANEL, line=color, radius=0.22)
    text(slide, x, y + h / 2 - 0.15, w, 0.3, label, size=9.5, bold=True,
         color=color, align=PP_ALIGN.CENTER)


def arrow(slide, x, y):
    a = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x), Inches(y),
                               Inches(0.15), Inches(0.19))
    a.rotation = 90
    a.fill.solid()
    a.fill.fore_color.rgb = MUTED
    a.line.fill.background()
    a.shadow.inherit = False


def slide_bg(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = box(s, 0, 0, 13.333, 7.5, BG)
    bg.shadow.inherit = False
    return s


# ---------------------------------------------------------------------- slides
def build():
    chart_progression(OUT / "chart_progression.png")
    chart_latency(OUT / "chart_latency.png")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # ------------------------------------------------------------- slide 1
    s = slide_bg(prs)
    box(s, 0, 0, 0.09, 7.5, ACCENT)
    text(s, 0.55, 0.40, 9.4, 0.55,
         "Video anomaly detection on a 4 GB laptop GPU",
         size=26, bold=True)
    text(s, 0.55, 1.04, 9.4, 0.34,
         [("Frozen vision-language encoder + two small trained heads  ·  11 classes  ·  "
           "20.3x real time on an RTX 3050 Laptop", {})],
         size=12.5, color=MUTED)

    stat(s, 10.25, 0.38, 2.55, "59.5 / 100", "PROJECTED MARKS · PUBLIC TEST")

    # pipeline
    text(s, 0.55, 1.72, 6, 0.3, "PIPELINE", size=10, bold=True, color=ACCENT)
    stages = [
        ("Drone / CCTV /\ndashcam video", "", MUTED),
        ("Sample 2 fps\ndownsize on decode", "OpenCV", MUTED),
        ("SigLIP-2 base\nFROZEN, fp16", "78 fps · 1.03 GB", ACCENT),
        ("4 s windows\nmean/max/std/drift", "2 s hop", MUTED),
        ("Two MLP heads\nwindow + clip", "trained, <1 min", GREEN),
        ("Hysteresis decoder\nrelative to baseline", "merge · gate", MUTED),
    ]
    x, w = 0.55, 1.92
    for i, (title, sub, col) in enumerate(stages):
        box(s, x, 2.12, w, 1.16, PANEL, line=col if col != MUTED else None, radius=0.1)
        text(s, x + 0.13, 2.27, w - 0.26, 0.6, title, size=10.2, bold=True,
             align=PP_ALIGN.CENTER, spacing=0.95)
        if sub:
            text(s, x + 0.13, 2.92, w - 0.26, 0.26, sub, size=8.4, color=col,
                 align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            arrow(s, x + w + 0.035, 2.61)
        x += w + 0.21

    text(s, 0.55, 3.52, 12.2, 0.3,
         [("Encode once, reuse forever.  ", {"bold": True, "color": ACCENT}),
          ("The encoder is frozen, so all 3,207 videos are embedded once into a cache. "
           "Retraining a head then takes under a minute, which is what made same-day "
           "iteration on windowing and thresholds possible.", {"color": MUTED})],
         size=11)

    # why not a VLM
    box(s, 0.55, 4.12, 6.05, 2.85, PANEL, radius=0.06)
    text(s, 0.85, 4.34, 5.45, 0.3, "WHY NOT RUN A VLM PER FRAME",
         size=10, bold=True, color=ACCENT)
    text(s, 0.85, 4.72, 5.45, 2.1,
         [("4 GB rules out fine-tuning a 2B VLM locally. So we froze one and trained "
           "only what sits on top.", {"size": 11.5, "color": TEXT}),
          ("The always-on stage is, in the end, a 12-way decision — orders of magnitude "
           "cheaper as a head over frozen embeddings than as generative decoding. The "
           "open-vocabulary semantics still come from the VLM; only the task-specific "
           "part is learned.",
           {"size": 11, "color": MUTED, "space_before": 9}),
          ("Explanations use nearest-neighbour retrieval over training descriptions in "
           "the same embedding space — one dot product per event, no generative model.",
           {"size": 11, "color": MUTED, "space_before": 9})],
         spacing=1.05)

    # right: latency + stats
    box(s, 6.85, 4.12, 5.95, 2.85, PANEL, radius=0.06)
    text(s, 7.15, 4.34, 5.4, 0.3, "SPEED, MEASURED END TO END",
         size=10, bold=True, color=ACCENT)
    s.shapes.add_picture(str(OUT / "chart_latency.png"), Inches(7.05), Inches(4.68),
                         height=Inches(1.62))
    text(s, 10.5, 4.78, 2.15, 1.5,
         [("3,391 s", {"size": 17, "bold": True, "color": TEXT}),
          ("of video", {"size": 9.5, "color": MUTED}),
          ("167 s", {"size": 17, "bold": True, "color": TEXT, "space_before": 7}),
          ("to process, incl. decode", {"size": 9.5, "color": MUTED})],
         spacing=1.0)
    text(s, 7.15, 6.42, 5.4, 0.4,
         [("Decode is the bottleneck, not the GPU. ",
           {"bold": True, "color": GREEN, "size": 10.5}),
          ("Encoding 2 fps costs 1.03 GB and leaves headroom for many streams.",
           {"color": MUTED, "size": 10.5})])

    # ------------------------------------------------------------- slide 2
    s = slide_bg(prs)
    box(s, 0, 0, 0.09, 7.5, ACCENT)
    text(s, 0.55, 0.42, 9, 0.45, "What moved the score, and what didn't",
         size=27, bold=True)
    text(s, 0.55, 1.0, 9, 0.3,
         "Every number below is the same 34-video public test set, scored locally "
         "against a re-implementation of the arena metric",
         size=12, color=MUTED)

    s.shapes.add_picture(str(OUT / "chart_progression.png"), Inches(0.55), Inches(1.5),
                         width=Inches(7.4))

    # levels
    y = 4.62
    text(s, 0.55, y, 7.4, 0.3, "FINAL, BY DIFFICULTY TIER",
         size=10, bold=True, color=ACCENT)
    for i, (lbl, val, marks, note) in enumerate([
            ("Difficulty 1", "0.750", "18.8 / 25", "anomaly acc 23 of 24"),
            ("Difficulty 2", "0.674", "23.6 / 35", "both normals silent"),
            ("Difficulty 3", "0.430", "17.2 / 40", "weakest, and worth most"),
    ]):
        bx = 0.55 + i * 2.52
        box(s, bx, y + 0.34, 2.32, 1.18, PANEL, radius=0.1)
        text(s, bx + 0.16, y + 0.45, 2.0, 0.24, lbl, size=9.5, color=MUTED)
        text(s, bx + 0.16, y + 0.68, 2.0, 0.36, val, size=19, bold=True)
        text(s, bx + 0.16, y + 1.06, 2.0, 0.24, marks + " marks", size=10,
             bold=True, color=ACCENT)
        text(s, bx + 0.16, y + 1.28, 2.05, 0.22, note, size=8, color=MUTED)

    text(s, 0.55, 6.35, 7.4, 0.85,
         [("Honest caveat.  ", {"bold": True, "color": ACCENT}),
          ("The decoder is tuned on 6 Difficulty-2 and 4 Difficulty-3 videos, so "
           "those thresholds are the least trustworthy part of the system, and ties "
           "break toward the least aggressive setting. Our scorer matches the "
           "published structure of the metric; the exact within-tier weighting "
           "isn't published, so it ranks configurations rather than predicting the "
           "leaderboard.", {"color": MUTED})],
         size=9.8, spacing=1.15)

    # findings column
    fx = 8.25
    text(s, fx, 1.5, 4.55, 0.3, "THREE THINGS THAT MATTERED",
         size=10, bold=True, color=ACCENT)
    findings = [
        ("Score against each video's own baseline",
         "The head, trained on short clips where the event fills the frame, "
         "reports a high flat pedestal on long footage — baselines ranged 0.55 to "
         "0.88, so no absolute threshold transfers. Thresholding on the rise above "
         "each video's own median found the events.",
         "L2 +0.12  ·  L3 +0.10", GREEN),
        ("Use each head for its own question",
         "The window head answers where; the clip head answers what. Letting the "
         "clip head classify each detected span, instead of pooling window votes, "
         "fixed whole videos at once.",
         "L3 0.296 → 0.414", GREEN),
        ("Merge gaps sized to real events",
         "Level-3 events run 45-125 s. A 5 s merge ceiling shattered one event into "
         "13 fragments — and only the best-overlapping fragment can ever match.",
         "T031 0.200 → 0.848", GREEN),
    ]
    yy = 1.9
    for title, body, gain, col in findings:
        box(s, fx, yy, 4.55, 1.42, PANEL, radius=0.08)
        text(s, fx + 0.2, yy + 0.13, 4.15, 0.26, title, size=11, bold=True)
        text(s, fx + 0.2, yy + 0.42, 4.15, 0.72, body, size=8.9, color=MUTED, spacing=1.08)
        text(s, fx + 0.2, yy + 1.14, 4.15, 0.24, gain, size=9, bold=True, color=col)
        yy += 1.54

    box(s, fx, 6.52, 4.55, 0.82, PANEL, radius=0.12)
    text(s, fx + 0.2, 6.63, 4.15, 0.22, "WHAT WE TRIED AND DROPPED",
         size=8.5, bold=True, color=RED)
    text(s, fx + 0.2, 6.87, 4.18, 0.4,
         "Larger SigLIP (segfaults at 4 GB) · zero-shot prompt fusion (0.29 alone, "
         "no lift) · one shared threshold for all tiers",
         size=8.2, color=MUTED, spacing=1.05)

    prs.save(OUT / "AHC_2slide_submission.pptx")
    print("wrote", OUT / "AHC_2slide_submission.pptx")


if __name__ == "__main__":
    build()
